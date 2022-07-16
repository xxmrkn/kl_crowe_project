import os
import numpy as np
import pandas as pd
import tqdm
import glob
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from itertools import product
import imageio
import cv2

def read_list_file(fpath):
    datalist = []
    if fpath.endswith('.txt'): #textfile
        with open(fpath, 'r' ,encoding="utf-8_sig", errors='ignore') as fr:
            datalist = fr.readlines()
        datalist = [dat.replace('\n', '') for dat in datalist]
    elif fpath.endswith('.csv'): #csv file
        df = pd.read_csv(fpath)
        print('Dataframe: ', df)
        #datalist = df.iloc[:,0].values.tolist()
        print(datalist)
    return datalist

def resize_img(img):
    scale_percent = 30 # percent of original size
    width = int(img.shape[1] * scale_percent / 100)
    height = int(img.shape[0] * scale_percent / 100)
    dim = (width, height)
    return cv2.resize(img, dim, interpolation = cv2.INTER_AREA)

def chunks(l, n):
    n = max(1, n)
    return (l[i:i+n] for i in range(0, len(l), n))

# _ROOT_IMGS = 'Z:/mazen/LowerLimbs/NII/20210827_NII_Pelvis/20210819_PelvisFemur_thumbs'
# _CSV = 'E:/Projects/Lower Limbs/Data/100_Osaka_Bones/patient_list_ver3.csv'
# df = pd.read_csv(_CSV, header=0, index_col=0)
# print(df)
_ROOT_IMG = 'C:/Users/masuda_m/code/20220511_DRR_with_Crowe_KL/DRR_AP'
# _ROOT_IMG = 'Z:/mazen/LowerLimbs/HipSegmetnation2/results/IwasaDataset/20211202_6_layer/Visualizations_Skin'
# _ROOT_LABEL = 'E:/Projects/Lower Limbs/NII/Pelvis/Visualizations/VolumeRendering_Revised'
# _ROOT_TGT_PPT = 'E:/Projects/Lower Limbs/NII/Pelvis/PPTs'
_ROOT_TGT_PPT = 'C:/Users/masuda_m/code' #+ '/PPT'
os.makedirs(_ROOT_TGT_PPT,exist_ok=True) 
# df = pd.read_csv(_CSV,header=0)

_ROWS = 5
_COLS = 10

_tops = [Cm(i) for i in np.arange(0,_ROWS*5,4.3)]
_tops_txt = [Cm(i) for i in np.arange(2.8,_ROWS*5,4.3)]
# print(len(_tops))
_ht = Cm(3)
_lefts_vol = [Cm(i) for i in np.arange(0,_COLS*3.1,3.1)]
_lefts_txt = [Cm(i) for i in np.arange(0,_COLS*3.1,3.1)]
# _lefts_drr = [Cm(i) for i in np.arange(0,_COLS*4,4)]
# _lefts_vol = [Cm(6), Cm(24)]
# _lefts_drr = [Cm(0), Cm(12)]

_MAX_SLIDES = 10
_MAX_CASES  = 3000
# _EXC_LIST_F = 'Z:/mazen/LowerLimbs/HipSegmetnation2/Data/Pass2/patient_list_phase3_diff.txt' 
_EXC_LIST_F = ''
_EXC_LIST = read_list_file(_EXC_LIST_F)
print(_EXC_LIST)
_IMG_EXT   = '_AP.png'

_LOG_FILE = os.path.join(_ROOT_TGT_PPT, 'log.txt')
_TMP_IMG = 'tmp_strata.png'
with open(_LOG_FILE, 'w') as f:
    f.writelines('Unprocessed list\n')

_EXTS = None
# imgs = glob.glob(os.path.join(_ROOT_IMG, '*_AP.png'))
# imgs.extend( glob.glob(os.path.join(_ROOT_IMG, '*_PA.png')))
# imgs.sort()
# for (_id, _path) in product(_EXC_LIST, imgs):
#     if _id in _path:
#         print(_id, _path)
#         imgs.remove(_path)
_CASE_LIST_F = 'C:/Users/masuda_m/code/path_VisionTransformer_Base16_2fold_10epoch_outlier.csv'
#_CASE_LIST_F = 'C:/Users/masuda_m/code/20220511_DRR_with_Crowe_KL/20220511_OsakaHip_TwoSide_KL_Crowe.csv'
# _CASE_LIST_F = 'Z:/mazen/LowerLimbs/HipSegmetnation2/Data/Phase2_Revised/caseid_list.txt'
# imgs = read_list_file(_CASE_LIST_F, field='KID')
_CASE_LIST = pd.read_csv(_CASE_LIST_F, header=0, index_col=0)
# _CASE_LIST.columns = ['img', 'grade']
print(f'_CASE_LIST{_CASE_LIST}')


# _EXTS = ['_right.png',
#          '_left.png']
#_CASE_LIST = _CASE_LIST.sort_values('Crowe', ascending=False)
imgs = _CASE_LIST.index.tolist()
# imgs.extend(glob.glob(os.path.join(_ROOT_IMG, '*R7*90.png')))

print(f'imgs{imgs}')
_IMG_CHUNKS = chunks(imgs,_MAX_CASES)
for h,_IMG_CHUNK in enumerate(_IMG_CHUNKS):
    prs = Presentation()
    prs.slide_width    = 11887200
    prs.slide_height   = 7786550
    title_slide_layout = prs.slide_layouts[5]
    blank_slide_layout = prs.slide_layouts[6]
    print('Image chunk #', h)
    _SLIDE_CHUNKS = chunks(_IMG_CHUNK,_ROWS*_COLS)
    pbar = tqdm.tqdm(_SLIDE_CHUNKS)
    for i,_IMAGE_LIST_CHUNK in enumerate(pbar):
        print('\tSlide chunk #', i)
        slide = prs.slides.add_slide(blank_slide_layout)
        _SUB_CHUNKS = chunks(_IMAGE_LIST_CHUNK, _COLS)
        for j, sub_chunk in enumerate(_SUB_CHUNKS):
            # print(sub_chunk)
            for _ID, img in enumerate(sub_chunk):
                try:
                    # _ID = '_'.join(os.path.basename(img_path).split('_')[0:3])
                    case_id = img
                    print(case_id)
                    # _ID = df.loc[k+1,'KID']
                    # print(_ID)
                    img_path = os.path.join(_ROOT_IMG, img)
                    #print(img_path)
                    if not os.path.exists(img_path):
                        img_path = os.path.join(_ROOT_IMG, img.lower())
                    # try:
                    #     if _EXTS:
                    #         imgs_f = []
                    #         for kk in [0]:
                    #             img_path = os.path.join(_ROOT_IMG, img )
                    #             img_vol = imageio.imread(img_path)
                    #             # img_vol = resize_img(img_vol)
                    #             # img_vol = np.append(img_vol, 255*np.ones(img_vol.shape[0],1), axis=1)
                    #         #     imgs_f.append(img_vol)
                    #         # img_vol = np.concatenate(imgs_f, axis=1)
                    #     else:
                    #         img_path = img
                    #         img_vol = imageio.imread(img_path)
                    #         img_vol = resize_img(img_vol)


                    # except:
                    #     with open(_LOG_FILE, 'a') as f:
                    #         f.writelines(img+'\n')
                    #     img_vol = None
                        # img_drr = None
                        # _COUNTER += 1
                    # imageio.imwrite('tmp_vol.png', img_vol)
                    # if img_drr is not None:
                    if img_path is not None:
                        # imageio.imwrite('tmp_drr.png', img_drr)
                        # imageio.imwrite('tmp_vol_2.png', img_vol)
                        # slide.shapes.add_picture('tmp_vol.png', _lefts_drr[k],_tops[j], height=_ht)
                        # slide.shapes.add_picture('tmp_vol.png', _lefts_img[_ID],_tops[j], height=_ht)
                        # slide.shapes.add_picture('tmp_vol_2.png', _lefts_vol[_ID],_tops[j], height=_ht)
                        slide.shapes.add_picture(img_path, _lefts_vol[_ID],_tops[j], height=_ht)
                    txBox = slide.shapes.add_textbox(_lefts_txt[_ID],_tops_txt[j],Cm(1),Cm(0.25))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    #print(case_id)
                    run.text = '%s\nActual: %s\nPred: %s\n' % (case_id, 
                                                        _CASE_LIST.loc[case_id, 'Actual'], 
                                                        _CASE_LIST.loc[case_id, 'Pred'])
                    font = run.font
                    font.name = 'Calibri'
                    font.size = Pt(8)
                    font.bold = True
                    font.color.theme_color = MSO_THEME_COLOR.ACCENT_5
                except Exception as e:
                    print(e)
                    txBox = slide.shapes.add_textbox(_lefts_txt[_ID],_tops_txt[j],Cm(1),Cm(0.25))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = '%s\nActual: %s\nPred: %s\n%s' % (case_id, 
                                                        _CASE_LIST.loc[case_id, 'Actual'], 
                                                        _CASE_LIST.loc[case_id, 'Pred'])
                    font = run.font
                    font.name = 'Calibri'
                    font.size = Pt(8)
                    font.bold = True
                    font.color.theme_color = MSO_THEME_COLOR.ACCENT_5
                    with open(_LOG_FILE, 'a') as f:
                        f.writelines('%s\n' % (img))

    prs.save(os.path.join(_ROOT_TGT_PPT, '20220701_Catalogue.pptx'))